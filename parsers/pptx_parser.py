from pathlib import Path

from pptx import Presentation


def parse_pptx(file_path: str | Path) -> str:
    """Extract text from slide shapes and return consolidated PPTX text."""
    presentation = Presentation(str(file_path))
    chunks: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(f"Slide {slide_number}: {shape.text.strip()}")

    return "\n".join(chunks).strip()